"""
文档结构化解析（Layer 0）
=========================
多源异构文档解析 → 目录抽取 → 目录关键词抽取 → 结构化输出。

支持文档类型：PDF / DOCX / TXT / PPTX
输出统一的 DocStructure JSON，供 Layer 1（doc_level1_extract）消费。

目录抽取策略（优先级从高到低）：
  1. DOCX heading 样式
  2. PPTX 幻灯片标题
  3. PDF/TXT 目录标记行（如 "..........."）
  4. LLM 逐页并发摘要 → 聚合生成目录（兜底）
"""

import os
import re
import sys
import json
import time
import hashlib
import traceback
from collections import defaultdict
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime
from threading import Lock

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

_V_DIR = os.path.dirname(os.path.abspath(__file__))
_PACKAGE_DIR = os.path.dirname(_V_DIR)
_EXTRACTION_DIR = os.path.dirname(_PACKAGE_DIR)
_SKILL_ROOT = os.path.dirname(_EXTRACTION_DIR)

try:
    from prompts import safe_parse_json_with_llm_repair
    from prompts_doc import doc_page_toc_summary_prompt, doc_toc_keywords_prompt
except ImportError:
    sys.path.insert(0, _SKILL_ROOT)
    sys.path.insert(0, _PACKAGE_DIR)
    from prompts import safe_parse_json_with_llm_repair
    from prompts_doc import doc_page_toc_summary_prompt, doc_toc_keywords_prompt


# ═══════════════════════════════════════════════════════════════════════════════
#  1. 多格式文档解析器
# ═══════════════════════════════════════════════════════════════════════════════

def _table_to_markdown(table_data: list[list]) -> str:
    """将二维表格数据转为 Markdown 表格字符串。"""
    if not table_data or not table_data[0]:
        return ""
    cleaned = []
    for row in table_data:
        cleaned.append([str(cell).strip().replace("\n", " ") if cell else "" for cell in row])

    col_count = max(len(r) for r in cleaned)
    for r in cleaned:
        while len(r) < col_count:
            r.append("")

    lines = []
    header = "| " + " | ".join(cleaned[0]) + " |"
    sep = "| " + " | ".join(["---"] * col_count) + " |"
    lines.append(header)
    lines.append(sep)
    for row in cleaned[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def parse_pdf(path: str) -> tuple[str, dict[int, str]]:
    """解析 PDF → (带页码标记的全文, {页码: 页面文本})。含表格识别。"""
    if pdfplumber is None:
        raise ImportError("pdfplumber 未安装，请运行: pip install pdfplumber")

    full_text = ""
    page_content: dict[int, str] = {}
    with pdfplumber.open(path) as pdf:
        for p, page in enumerate(pdf.pages, start=1):
            page_text = page.extract_text() or ""

            tables = page.extract_tables() or []
            if tables:
                table_md_parts = []
                for tbl in tables:
                    md = _table_to_markdown(tbl)
                    if md:
                        table_md_parts.append(md)
                if table_md_parts:
                    page_text += "\n\n" + "\n\n".join(table_md_parts)

            full_text += f"page:{p}\n\n{page_text}\n" + "—" * 50 + "\n"
            page_content[p] = page_text
    return full_text, page_content


def parse_docx(path: str) -> tuple[str, dict[int, str]]:
    """解析 DOCX → (全文, {段落序号: 段落文本})。含表格识别，按文档内元素顺序输出。"""
    if DocxDocument is None:
        raise ImportError("python-docx 未安装，请运行: pip install python-docx")

    doc = DocxDocument(path)
    full_text = ""
    segment_content: dict[int, str] = {}
    seg_idx = 0

    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            from docx.text.paragraph import Paragraph
            para = Paragraph(element, doc)
            text = para.text.strip()
            if not text:
                continue
            seg_idx += 1
            full_text += f"segment:{seg_idx}\n\n{text}\n" + "—" * 50 + "\n"
            segment_content[seg_idx] = text

        elif tag == "tbl":
            from docx.table import Table
            table = Table(element, doc)
            rows = []
            for row in table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            md = _table_to_markdown(rows)
            if md:
                seg_idx += 1
                full_text += f"segment:{seg_idx}\n\n{md}\n" + "—" * 50 + "\n"
                segment_content[seg_idx] = md

    return full_text, segment_content


def parse_docx_with_styles(path: str) -> list[dict]:
    """解析 DOCX 并保留段落样式（含表格），返回 [{idx, text, style, level}]。"""
    if DocxDocument is None:
        raise ImportError("python-docx 未安装，请运行: pip install python-docx")

    doc = DocxDocument(path)
    paragraphs = []
    seg_idx = 0

    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            from docx.text.paragraph import Paragraph
            para = Paragraph(element, doc)
            text = para.text.strip()
            if not text:
                continue
            seg_idx += 1
            style_name = para.style.name if para.style else ""
            level = -1
            if style_name.startswith("Heading"):
                try:
                    level = int(style_name.replace("Heading", "").strip())
                except ValueError:
                    level = 0
            paragraphs.append({
                "idx": seg_idx,
                "text": text,
                "style": style_name,
                "level": level,
            })

        elif tag == "tbl":
            from docx.table import Table
            table = Table(element, doc)
            rows = []
            for row in table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            md = _table_to_markdown(rows)
            if md:
                seg_idx += 1
                paragraphs.append({
                    "idx": seg_idx,
                    "text": md,
                    "style": "Table",
                    "level": -1,
                })

    return paragraphs


def parse_txt(path: str) -> tuple[str, dict[int, str]]:
    """解析 TXT → (全文, {段落序号: 段落文本})。按双换行分段。"""
    with open(path, "r", encoding="utf-8") as f:
        raw = f.read()

    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", raw) if p.strip()]
    full_text = ""
    segment_content: dict[int, str] = {}
    for i, para in enumerate(paragraphs, start=1):
        full_text += f"segment:{i}\n\n{para}\n" + "—" * 50 + "\n"
        segment_content[i] = para
    return full_text, segment_content


def parse_pptx(path: str) -> tuple[str, dict[int, str]]:
    """解析 PPTX → (全文, {幻灯片序号: 幻灯片文本})。含表格识别。"""
    if Presentation is None:
        raise ImportError("python-pptx 未安装，请运行: pip install python-pptx")

    prs = Presentation(path)
    full_text = ""
    slide_content: dict[int, str] = {}
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        parts.append(t)
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                md = _table_to_markdown(rows)
                if md:
                    parts.append(md)
        page_text = "\n".join(parts)
        full_text += f"slide:{i}\n\n{page_text}\n" + "—" * 50 + "\n"
        slide_content[i] = page_text
    return full_text, slide_content


_PARSER_MAP = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".txt": parse_txt,
    ".pptx": parse_pptx,
}


def merge_segments_by_length(
    segment_content: dict[int, str],
    min_chars: int = 500,
    max_chars: int = 2000,
) -> tuple[dict[int, str], str]:
    """
    按字数窗口合并过短的段落，确保每个 segment 在 [min_chars, max_chars] 区间。

    规则：
    - 依次累积段落到 buffer
    - 若加入当前段落后 buffer 超过 max_chars → 先把已有 buffer 输出，
      当前段落作为新 buffer 的起点（若单段本身 > max_chars，直接独立输出）
    - 若 buffer 累积到 >= min_chars → 输出
    - 末尾剩余不足 min_chars → 尝试合并到上一个 segment（不超 max_chars 时），否则独立输出

    Returns
    -------
    merged_segments : {新序号: 合并后文本}
    full_text       : 重建的带 segment 标记全文
    """
    sorted_keys = sorted(segment_content.keys())
    if not sorted_keys:
        return {}, ""

    merged: dict[int, str] = {}
    seg_id = 1
    buffer_parts: list[str] = []
    buffer_len = 0

    for idx in sorted_keys:
        text = segment_content[idx]
        text_len = len(text)

        if buffer_parts and buffer_len + text_len > max_chars:
            merged[seg_id] = "\n\n".join(buffer_parts)
            seg_id += 1
            buffer_parts = []
            buffer_len = 0

        buffer_parts.append(text)
        buffer_len += text_len

        if text_len > max_chars and len(buffer_parts) == 1:
            merged[seg_id] = "\n\n".join(buffer_parts)
            seg_id += 1
            buffer_parts = []
            buffer_len = 0
        elif buffer_len >= min_chars:
            merged[seg_id] = "\n\n".join(buffer_parts)
            seg_id += 1
            buffer_parts = []
            buffer_len = 0

    if buffer_parts:
        leftover = "\n\n".join(buffer_parts)
        if merged and buffer_len < min_chars:
            last_key = max(merged.keys())
            if len(merged[last_key]) + len(leftover) + 2 <= max_chars:
                merged[last_key] += "\n\n" + leftover
            else:
                merged[seg_id] = leftover
        else:
            merged[seg_id] = leftover

    full_text = ""
    for sid in sorted(merged.keys()):
        full_text += f"segment:{sid}\n\n{merged[sid]}\n" + "—" * 50 + "\n"

    return merged, full_text


def parse_document(path: str) -> tuple[str, dict[int, str], str]:
    """
    统一入口：根据扩展名自动选择解析器。

    Returns
    -------
    full_text       : 带段落/页码标记的全文
    segment_content : {序号: 文本}
    file_type       : pdf / docx / txt / pptx
    """
    ext = os.path.splitext(path)[1].lower()
    if ext not in _PARSER_MAP:
        raise ValueError(f"不支持的文档类型: {ext}（支持 {list(_PARSER_MAP.keys())}）")
    full_text, segments = _PARSER_MAP[ext](path)
    return full_text, segments, ext.lstrip(".")


# ═══════════════════════════════════════════════════════════════════════════════
#  2. 目录抽取：规则优先 → LLM 逐页并发摘要兜底
# ═══════════════════════════════════════════════════════════════════════════════

def _extract_toc_by_marker(text: str, segment_content: dict, toc_marker: str = "...........") -> list[dict]:
    """
    规则方式：通过连续省略号等标记识别目录行（适用于结构化 PDF）。
    返回 [{title, level, start_seg, end_seg}]。
    """
    menu = [line for line in text.split("\n") if toc_marker in line]
    if not menu:
        return []

    total_segs = max(segment_content.keys()) if segment_content else 1
    toc_items = []

    for i in range(len(menu)):
        title = menu[i].split("..")[0].strip()
        try:
            start = int(re.sub(r"\D", "", menu[i][-10:]))
        except (ValueError, IndexError):
            continue
        if i + 1 < len(menu):
            try:
                end = int(re.sub(r"\D", "", menu[i + 1][-10:]))
            except (ValueError, IndexError):
                end = total_segs
        else:
            end = total_segs

        level = _infer_heading_level(title)
        toc_items.append({
            "title": title,
            "level": level,
            "start_seg": start,
            "end_seg": end,
        })

    return toc_items


def _extract_toc_from_docx_styles(path: str) -> list[dict]:
    """从 DOCX heading 样式直接提取目录结构。"""
    styled_paras = parse_docx_with_styles(path)
    toc_items = []
    heading_paras = [p for p in styled_paras if p["level"] >= 0]

    for i, hp in enumerate(heading_paras):
        start_idx = hp["idx"]
        if i + 1 < len(heading_paras):
            end_idx = heading_paras[i + 1]["idx"] - 1
        else:
            end_idx = styled_paras[-1]["idx"] if styled_paras else start_idx

        toc_items.append({
            "title": hp["text"],
            "level": hp["level"],
            "start_seg": start_idx,
            "end_seg": end_idx,
        })

    return toc_items


def _extract_toc_from_pptx(path: str) -> list[dict]:
    """从 PPTX 幻灯片标题提取目录结构（每张幻灯片视为一个段落）。"""
    if Presentation is None:
        return []
    prs = Presentation(path)
    toc_items = []
    for i, slide in enumerate(prs.slides, start=1):
        title_text = ""
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
        if not title_text:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        title_text = t[:80]
                        break
        if title_text:
            toc_items.append({
                "title": title_text,
                "level": 1,
                "start_seg": i,
                "end_seg": i,
            })
    return toc_items


# ─── LLM 逐页并发摘要 → 聚合目录 ────────────────────────────────────────────

_toc_llm_lock = Lock()


def _summarize_single_page(
    page_num: int,
    page_text: str,
    llm_func,
    max_retries: int = 3,
) -> list[dict]:
    """
    对单页内容调用 LLM，生成该页的目录标题列表。
    返回 [{"title": ..., "level": ..., "page": page_num}]。
    """
    if not page_text.strip():
        return []

    for attempt in range(max_retries):
        try:
            prompt = doc_page_toc_summary_prompt(page_num, page_text)
            response = llm_func(prompt)
            raw = response.get("content", "") if isinstance(response, dict) else str(response)
            parsed = safe_parse_json_with_llm_repair(raw, llm_func=llm_func)

            results = []
            for t in parsed.get("titles", []):
                title = t.get("title", "").strip()
                if title:
                    results.append({
                        "title": title,
                        "level": int(t.get("level", 1)),
                        "page": page_num,
                    })
            return results

        except Exception:
            if attempt < max_retries - 1:
                wait = min(2 ** attempt, 15)
                time.sleep(wait)
            else:
                print(f"  [TOC-LLM] 第 {page_num} 页摘要失败（{max_retries}次重试耗尽）: "
                      f"{traceback.format_exc()[:120]}")
                return []


def _aggregate_page_titles_to_toc(
    page_titles: list[dict],
    total_segs: int,
) -> list[dict]:
    """
    将逐页 LLM 摘要结果聚合成完整目录。

    逻辑：
    - 按页码排序后，相邻页面中相同标题合并为一个条目
    - 每个条目的 start_seg = 首次出现的页码，end_seg = 最后出现的页码
    - 对无标题覆盖的页码间隙，扩展前一个条目的 end_seg
    """
    if not page_titles:
        return [{"title": "全文", "level": 1, "start_seg": 1, "end_seg": total_segs}]

    sorted_titles = sorted(page_titles, key=lambda x: (x["page"], x["level"]))

    toc_items = []
    prev_title = None
    for entry in sorted_titles:
        if prev_title and entry["title"] == prev_title["title"]:
            prev_title["end_seg"] = entry["page"]
        else:
            if prev_title:
                toc_items.append(prev_title)
            prev_title = {
                "title": entry["title"],
                "level": entry["level"],
                "start_seg": entry["page"],
                "end_seg": entry["page"],
            }
    if prev_title:
        toc_items.append(prev_title)

    for i in range(len(toc_items) - 1):
        if toc_items[i]["end_seg"] < toc_items[i + 1]["start_seg"] - 1:
            toc_items[i]["end_seg"] = toc_items[i + 1]["start_seg"] - 1

    if toc_items and toc_items[-1]["end_seg"] < total_segs:
        toc_items[-1]["end_seg"] = total_segs

    return toc_items


def _extract_toc_by_llm(
    segment_content: dict,
    llm_func,
    max_workers: int = 8,
    max_retries: int = 3,
) -> list[dict]:
    """
    LLM 兜底：按页/段并发调用 LLM 生成逐页目录标题，再聚合成完整目录。
    """
    total_segs = max(segment_content.keys()) if segment_content else 1
    seg_keys = sorted(segment_content.keys())

    print(f"  [TOC-LLM] 逐页并发摘要，共 {len(seg_keys)} 页/段，并发数: {max_workers}")

    all_page_titles: list[dict] = []

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_to_page = {
            executor.submit(
                _summarize_single_page,
                page_num, segment_content[page_num], llm_func, max_retries,
            ): page_num
            for page_num in seg_keys
        }
        completed = 0
        for future in as_completed(future_to_page):
            page_num = future_to_page[future]
            try:
                titles = future.result()
                all_page_titles.extend(titles)
                completed += 1
                if completed % 10 == 0 or completed == len(seg_keys):
                    print(f"  [TOC-LLM] 进度: {completed}/{len(seg_keys)}")
            except Exception as e:
                print(f"  [TOC-LLM] 第 {page_num} 页异常: {e}")

    print(f"  [TOC-LLM] 收集到 {len(all_page_titles)} 个页级标题，开始聚合...")
    toc_items = _aggregate_page_titles_to_toc(all_page_titles, total_segs)
    return toc_items


def _infer_heading_level(title: str) -> int:
    """根据标题格式推断层级（第X章=1, X.Y=2, X.Y.Z=3 ...）。"""
    if re.match(r"^第\s*\d+\s*[章篇部]", title):
        return 1
    m = re.match(r"^(\d+(?:\.\d+)*)", title)
    if m:
        return m.group(1).count(".") + 1
    return 1


def extract_toc(
    doc_path: str,
    full_text: str,
    segment_content: dict,
    file_type: str,
    llm_func=None,
    toc_marker: str = "...........",
    llm_toc_workers: int = 8,
    force_llm_toc: bool = False,
) -> list[dict]:
    """
    统一目录抽取入口，策略优先级：
    1. DOCX heading 样式
    2. PPTX 幻灯片标题
    3. PDF/TXT 目录标记行
    4. LLM 逐页并发摘要 → 聚合目录（兜底）

    当 force_llm_toc=True 时，跳过所有规则方式，直接走 LLM 逐页摘要。
    """
    if force_llm_toc:
        if llm_func is None:
            raise ValueError("force_llm_toc=True 但未提供 llm_func")
        print("  [TOC] force_llm_toc=True，跳过规则方式，直接使用 LLM 逐页并发摘要...")
        toc = _extract_toc_by_llm(segment_content, llm_func, max_workers=llm_toc_workers)
        if toc:
            print(f"  [TOC] LLM 逐页摘要聚合: {len(toc)} 条")
            return toc
        total_segs = max(segment_content.keys()) if segment_content else 1
        return [{"title": "全文", "level": 1, "start_seg": 1, "end_seg": total_segs}]

    toc = []

    if file_type == "docx":
        toc = _extract_toc_from_docx_styles(doc_path)
        if toc:
            print(f"  [TOC] DOCX heading 样式提取: {len(toc)} 条")
            return toc

    if file_type == "pptx":
        toc = _extract_toc_from_pptx(doc_path)
        if toc:
            print(f"  [TOC] PPTX 幻灯片标题提取: {len(toc)} 条")
            return toc

    toc = _extract_toc_by_marker(full_text, segment_content, toc_marker)
    if toc:
        print(f"  [TOC] 目录标记行提取: {len(toc)} 条")
        return toc

    if llm_func is not None:
        print("  [TOC] 无显式目录标记，启动 LLM 逐页并发摘要...")
        toc = _extract_toc_by_llm(segment_content, llm_func, max_workers=llm_toc_workers)
        if toc:
            print(f"  [TOC] LLM 逐页摘要聚合: {len(toc)} 条")
            return toc

    print("  [TOC] 警告：未能提取任何目录结构，将整篇文档视为单一章节。")
    total_segs = max(segment_content.keys()) if segment_content else 1
    return [{"title": "全文", "level": 1, "start_seg": 1, "end_seg": total_segs}]


# ═══════════════════════════════════════════════════════════════════════════════
#  3. 目录关键词抽取
# ═══════════════════════════════════════════════════════════════════════════════

def _process_keyword_batch(
    batch_start: int,
    batch: list[dict],
    segment_content: dict,
    llm_func,
) -> list[tuple[int, list[str]]]:
    """
    处理单个关键词批次：构造 prompt → 调用 LLM → 返回 [(全局索引, 关键词列表)]。
    """
    entries = []
    for item in batch:
        seg_range = range(item["start_seg"], item["end_seg"] + 1)
        snippet = "\n".join(
            segment_content.get(s, "")[:200] for s in seg_range
        )[:600]
        entries.append({"title": item["title"], "snippet": snippet})

    prompt = doc_toc_keywords_prompt(entries)
    response = llm_func(prompt)
    raw = response.get("content", "") if isinstance(response, dict) else str(response)
    parsed = safe_parse_json_with_llm_repair(raw, llm_func=llm_func)
    kw_list = parsed.get("keywords", [])

    results = []
    for i, kw in enumerate(kw_list):
        global_idx = batch_start + i
        results.append((global_idx, kw if isinstance(kw, list) else [kw]))
    return results


def extract_toc_keywords(
    toc_items: list[dict],
    segment_content: dict,
    llm_func=None,
    batch_size: int = 1,
    max_workers: int = 8,
) -> list[dict]:
    """
    为每个目录条目提取关键词（并发）。
    - 若提供 llm_func，按 batch_size 分批并发调用 LLM 提取；
    - batch_size=1（默认）时每个条目独立调用，上下文最短、提取质量最高；
    - 否则基于标题本身做简单关键词切分。

    直接在 toc_items 上增加 "keywords" 字段并返回。
    """
    if llm_func is None:
        for item in toc_items:
            item["keywords"] = _rule_based_keywords(item["title"])
        return toc_items

    batches = []
    for batch_start in range(0, len(toc_items), batch_size):
        batch = toc_items[batch_start: batch_start + batch_size]
        batches.append((batch_start, batch))

    print(f"  [Keywords] 共 {len(toc_items)} 条目录，分 {len(batches)} 批并发提取关键词"
          f"（batch_size={batch_size}, workers={max_workers}）")

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_to_batch = {
            executor.submit(
                _process_keyword_batch,
                batch_start, batch, segment_content, llm_func,
            ): batch_start
            for batch_start, batch in batches
        }
        for future in as_completed(future_to_batch):
            batch_start = future_to_batch[future]
            try:
                results = future.result()
                for global_idx, kw in results:
                    if global_idx < len(toc_items):
                        toc_items[global_idx]["keywords"] = kw
            except Exception as e:
                print(f"  [Keywords] LLM 关键词提取失败（batch {batch_start}）: {e}")
                batch_end = min(batch_start + batch_size, len(toc_items))
                for i in range(batch_start, batch_end):
                    toc_items[i].setdefault("keywords", _rule_based_keywords(toc_items[i]["title"]))

    for item in toc_items:
        item.setdefault("keywords", _rule_based_keywords(item["title"]))

    return toc_items


def _rule_based_keywords(title: str) -> list[str]:
    """简单规则：按标点和常见停用词切分标题。"""
    cleaned = re.sub(r"[第\d+章篇部节条\.\s\-—：:（）\(\)]", " ", title)
    tokens = [t.strip() for t in cleaned.split() if len(t.strip()) >= 2]
    return tokens if tokens else [title.strip()]


# ═══════════════════════════════════════════════════════════════════════════════
#  4. 段落映射（段落 → 目录节）
# ═══════════════════════════════════════════════════════════════════════════════

def build_paragraphs(segment_content: dict[int, str], toc_items: list[dict]) -> list[dict]:
    """
    将原始段落映射到所属目录节点，构建段落列表。
    返回 [{idx, text, toc_section, toc_level}]。
    """
    seg_to_toc: dict[int, dict] = {}
    for item in toc_items:
        for s in range(item["start_seg"], item["end_seg"] + 1):
            if s not in seg_to_toc:
                seg_to_toc[s] = item

    paragraphs = []
    for idx in sorted(segment_content.keys()):
        text = segment_content[idx]
        if not text.strip():
            continue
        toc_info = seg_to_toc.get(idx, toc_items[0] if toc_items else {"title": "未知", "level": 0})
        paragraphs.append({
            "idx": idx,
            "text": text,
            "toc_section": toc_info["title"],
            "toc_level": toc_info.get("level", 0),
        })
    return paragraphs


# ═══════════════════════════════════════════════════════════════════════════════
#  5. 主入口：生成 DocStructure JSON
# ═══════════════════════════════════════════════════════════════════════════════

def run_doc_structure_parse(
    doc_path: str,
    llm_func=None,
    output_file: str = None,
    toc_marker: str = "...........",
    extract_keywords: bool = True,
    llm_toc_workers: int = 8,
    force_llm_toc: bool = False,
    min_seg_chars: int = 500,
    max_seg_chars: int = 2000,
) -> dict:
    """
    Layer 0 主入口：从原始文档生成结构化中间表示。

    Parameters
    ----------
    doc_path          : 文档文件路径（PDF/DOCX/TXT/PPTX）
    llm_func          : LLM 调用函数（用于 TOC 推断、关键词提取；为 None 则仅用规则）
    output_file       : 输出 JSON 路径（为 None 则不写文件）
    toc_marker        : PDF 目录行标记
    extract_keywords  : 是否提取目录关键词
    llm_toc_workers   : LLM 逐页摘要的并发线程数
    force_llm_toc     : 强制使用 LLM 逐页摘要生成目录（跳过所有规则方式）
    min_seg_chars     : 段落合并下限（字数不足时向后拼接），0 则不合并
    max_seg_chars     : 段落合并上限（超过时独立切出）

    Returns
    -------
    DocStructure dict
    """
    print(f"[Layer-0] 开始文档结构化解析: {os.path.basename(doc_path)}")

    # Step 1: 解析文档
    print("  [1/5] 解析文档原始内容...")
    full_text, segment_content, file_type = parse_document(doc_path)
    raw_seg_count = len(segment_content)
    print(f"  文档类型: {file_type}，原始段落/页数: {raw_seg_count}")

    # Step 2: 按字数窗口合并段落
    if min_seg_chars > 0:
        print(f"  [2/5] 按字数窗口合并段落（min={min_seg_chars}, max={max_seg_chars}）...")
        segment_content, full_text = merge_segments_by_length(
            segment_content, min_chars=min_seg_chars, max_chars=max_seg_chars,
        )
        merged_seg_count = len(segment_content)
        print(f"  合并结果: {raw_seg_count} 段 → {merged_seg_count} 段")
    else:
        print("  [2/5] 跳过段落合并（min_seg_chars=0）")

    # Step 3: 目录抽取
    print("  [3/5] 抽取目录结构...")
    toc_items = extract_toc(
        doc_path, full_text, segment_content, file_type,
        llm_func, toc_marker, llm_toc_workers, force_llm_toc,
    )

    # Step 4: 关键词抽取
    if extract_keywords:
        print("  [4/5] 抽取目录关键词...")
        toc_items = extract_toc_keywords(toc_items, segment_content, llm_func)
    else:
        print("  [4/5] 跳过关键词抽取")
        for item in toc_items:
            item.setdefault("keywords", _rule_based_keywords(item["title"]))

    # Step 5: 段落映射
    print("  [5/5] 构建段落-目录映射...")
    paragraphs = build_paragraphs(segment_content, toc_items)

    # 组装输出结构
    doc_structure = {
        "document_meta": {
            "source_file": os.path.abspath(doc_path),
            "file_name": os.path.basename(doc_path),
            "file_type": file_type,
            "total_segments": len(segment_content),
            "parse_timestamp": datetime.now().isoformat(),
            "content_hash": hashlib.md5(full_text.encode("utf-8")).hexdigest(),
        },
        "toc": toc_items,
        "paragraphs": [
            {
                "idx": p["idx"],
                "text": p["text"],
                "toc_section": p["toc_section"],
                "toc_level": p["toc_level"],
            }
            for p in paragraphs
        ],
        "parse_config": {
            "toc_marker": toc_marker,
            "min_seg_chars": min_seg_chars,
            "max_seg_chars": max_seg_chars,
        },
    }

    if output_file:
        os.makedirs(os.path.dirname(output_file) or ".", exist_ok=True)
        with open(output_file, "w", encoding="utf-8") as f:
            json.dump(doc_structure, f, ensure_ascii=False, indent=2)
        print(f"[Layer-0] 结构化结果已保存: {output_file}")

    _print_summary(doc_structure)
    return doc_structure


def _print_summary(ds: dict):
    """打印 DocStructure 摘要。"""
    meta = ds["document_meta"]
    toc = ds["toc"]
    paras = ds["paragraphs"]

    print(f"\n{'═' * 60}")
    print(f"  文档: {meta['file_name']} ({meta['file_type']})")
    print(f"  总段落/页: {meta['total_segments']}  |  目录条目: {len(toc)}  |  映射段落: {len(paras)}")

    print(f"  目录层级分布: ", end="")
    level_counts = defaultdict(int)
    for t in toc:
        level_counts[t["level"]] += 1
    for lv in sorted(level_counts):
        print(f"L{lv}={level_counts[lv]} ", end="")

    section_counts = defaultdict(int)
    for p in paras:
        section_counts[p["toc_section"]] += 1
    avg_paras = sum(section_counts.values()) / max(len(section_counts), 1)
    print(f"\n  每目录节平均段落数: {avg_paras:.1f}")
    print(f"{'═' * 60}\n")


# ═══════════════════════════════════════════════════════════════════════════════
#  独立运行入口：支持指定文件 or 扫描 input 文件夹
# ═══════════════════════════════════════════════════════════════════════════════

_SUPPORTED_DOC_EXTS = set(_PARSER_MAP.keys())


if __name__ == "__main__":
    import argparse

    sys.path.insert(0, _SKILL_ROOT)
    from llm_client import chat
    sys.path.insert(0, _EXTRACTION_DIR)
    from utils import get_source_stem

    parser = argparse.ArgumentParser(
        description="文档结构化解析（Layer 0）"
    )
    parser.add_argument(
        "--files", "-f", nargs="+", default=None,
        help="指定要处理的文档文件路径（支持多个）；不指定则处理 input 目录下所有文件",
    )
    args = parser.parse_args()

    input_dir = os.path.join(_PACKAGE_DIR, "input")
    output_dir = os.path.join(_PACKAGE_DIR, "output")
    os.makedirs(output_dir, exist_ok=True)

    if args.files:
        doc_files = []
        for fp in args.files:
            fp = os.path.abspath(fp)
            if not os.path.isfile(fp):
                print(f"[警告] 文件不存在，已跳过: {fp}")
                continue
            if os.path.splitext(fp)[1].lower() not in _SUPPORTED_DOC_EXTS:
                print(f"[警告] 不支持的文件类型，已跳过: {fp}（支持: {_SUPPORTED_DOC_EXTS}）")
                continue
            doc_files.append(fp)
        if not doc_files:
            raise FileNotFoundError("指定的文件中没有可处理的有效文件")
        mode_desc = "指定文件模式"
    else:
        doc_files = sorted([
            os.path.join(input_dir, f)
            for f in os.listdir(input_dir)
            if os.path.splitext(f)[1].lower() in _SUPPORTED_DOC_EXTS
        ])
        if not doc_files:
            raise FileNotFoundError(
                f"input 目录中未找到支持的文档文件（{_SUPPORTED_DOC_EXTS}）：{input_dir}"
            )
        mode_desc = "全量扫描模式"

    print("=" * 60)
    print(f"[doc_structure_parse] {mode_desc}，共 {len(doc_files)} 个源文档（仅结构化解析）")
    print("=" * 60)
    for i, fp in enumerate(doc_files, 1):
        print(f"  {i}. {os.path.basename(fp)}")

    for idx, doc_path in enumerate(doc_files, 1):
        print(f"\n{'═' * 60}")
        print(f"  [{idx}/{len(doc_files)}] 处理: {os.path.basename(doc_path)}")
        print(f"{'═' * 60}")

        source_stem = get_source_stem(doc_path)
        output_file = os.path.join(output_dir, f"{source_stem}_structure.json")

        ds = run_doc_structure_parse(
            doc_path=doc_path,
            llm_func=chat,
            output_file=output_file,
        )

        print(f"\n  目录结构预览（前 15 条）:")
        for item in ds["toc"][:15]:
            indent = "  " * item["level"]
            kw = ", ".join(item.get("keywords", []))
            print(f"    {indent}[L{item['level']}] {item['title']}  "
                  f"(seg {item['start_seg']}~{item['end_seg']})  关键词: {kw}")

    print(f"\n{'═' * 60}")
    print(f"[doc_structure_parse] 全部 {len(doc_files)} 个文档处理完成！")
    print(f"{'═' * 60}")
